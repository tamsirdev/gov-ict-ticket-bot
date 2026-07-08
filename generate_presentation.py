"""
Generate a PowerPoint presentation for the Government ICT Ticket Bot.

Usage:
    pip install python-pptx
    python generate_presentation.py
"""

from datetime import datetime

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
except ImportError:
    msg = (
        "python-pptx is required. Install it with: pip install python-pptx"
    )
    raise ImportError(msg)


DARK_BG = RGBColor(0x1A, 0x20, 0x2C)
ACCENT_BLUE = RGBColor(0x63, 0xB3, 0xED)
LIGHT_BLUE = RGBColor(0x90, 0xCD, 0xF4)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0xA0, 0xAE, 0xC0)
DARK_CARD = RGBColor(0x2D, 0x37, 0x48)
GREEN = RGBColor(0x27, 0x67, 0x49)
RED = RGBColor(0x9B, 0x2C, 0x2C)


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height


def add_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def text_box(slide, left, top, width, height, text, size=18, bold=False,
             color=WHITE, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_bullet_slide(slide, left, top, width, height, items, size=16,
                     color=WHITE):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(6)
    return tf


def title_slide(prs, title, subtitle, date_str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    text_box(slide, 1, 2, 11, 1.5, title, size=44, bold=True,
             color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)
    text_box(slide, 1, 3.5, 11, 1, subtitle, size=22,
             color=LIGHT_BLUE, alignment=PP_ALIGN.CENTER)
    text_box(slide, 1, 5.5, 11, 0.5, date_str, size=14,
             color=GRAY, alignment=PP_ALIGN.CENTER)


def section_slide(prs, number, title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    text_box(slide, 1, 1, 2, 0.5, f"0{number}", size=60, bold=True,
             color=ACCENT_BLUE)
    text_box(slide, 1, 2.5, 11, 1, title, size=36, bold=True, color=WHITE)
    line = slide.shapes.add_shape(
        1, Inches(1), Inches(3.5), Inches(3), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_BLUE
    line.line.fill.background()


def content_slide(prs, title, items, two_col=False):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    text_box(slide, 0.8, 0.4, 11, 0.7, title, size=28, bold=True,
             color=ACCENT_BLUE)
    line = slide.shapes.add_shape(
        1, Inches(0.8), Inches(1.1), Inches(2), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_BLUE
    line.line.fill.background()

    if two_col:
        mid = len(items) // 2 + len(items) % 2
        add_bullet_slide(slide, 0.8, 1.5, 5.5, 5, items[:mid], size=16)
        add_bullet_slide(slide, 6.8, 1.5, 5.5, 5, items[mid:], size=16)
    else:
        add_bullet_slide(slide, 0.8, 1.5, 11, 5, items, size=16)
    return slide


# ── Title ──────────────────────────────────────────────────────────────
title_slide(prs,
    "Government ICT Service Ticket Bot",
    "Automated Helpdesk with AI Triage & DevOps Pipeline",
    datetime.now().strftime("%B %Y"))

# ── Problem ────────────────────────────────────────────────────────────
section_slide(prs, 1, "The Problem")
content_slide(prs, "ICT Helpdesk Challenges", [
    "• ICT helpdesks in government are overwhelmed with manual ticket handling",
    "• Citizens struggle to find the right channel to report issues",
    "• ICT officers waste time categorizing and prioritizing tickets manually",
    "• No automated tracking — tickets get lost in emails and phone calls",
    "",
    '  "A citizen\'s printer has been down for two weeks',
    '   because the ticket was never logged."',
])

# ── Solution ───────────────────────────────────────────────────────────
section_slide(prs, 2, "The Solution")
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
text_box(slide, 0.8, 0.4, 11, 0.7, "For Citizens & ICT Officers", size=28,
         bold=True, color=ACCENT_BLUE)
line = slide.shapes.add_shape(
    1, Inches(0.8), Inches(1.1), Inches(2), Inches(0.03))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT_BLUE
line.line.fill.background()
add_bullet_slide(slide, 0.8, 1.5, 5.5, 5, [
    "For Citizens",
    "  • Submit issues via web browser or WhatsApp",
    "  • Get instant ticket ID & confirmation",
    "  • Check status anytime",
    "",
    "For ICT Officers",
    "  • Admin dashboard with all tickets in one place",
    "  • Auto-categorized (Hardware, Software, Network, Account)",
    "  • Priority-assigned — Critical issues surface first",
], size=16)

# ── Features ───────────────────────────────────────────────────────────
section_slide(prs, 3, "Features")

content_slide(prs, "Free vs Paid", [
    "Free Mode (no API keys):",
    "  • Web form submission",
    "  • Rule-based keyword triage",
    "  • Admin dashboard",
    "  • Docker + CI/CD",
    "  • $0 cost",
    "",
    "Paid Upgrades (when API key set):",
    "  • OpenAI GPT-3.5 triage (smarter categorization)",
    "  • WhatsApp integration via Twilio",
])

# ── Architecture ───────────────────────────────────────────────────────
section_slide(prs, 4, "System Architecture")
content_slide(prs, "Components", [
    "• Flask Application (bot.py) — HTTP routes, web UI, admin, webhook",
    "• Triage Engine (ai_triage.py) — Free rules + optional OpenAI",
    "• Database (database.py) — SQLite CRUD operations",
    "• Templates (templates/) — Jinja2 HTML pages",
    "• Tests (tests/) — 27 pytest tests across 3 test files",
])

# ── Tech Stack ─────────────────────────────────────────────────────────
section_slide(prs, 5, "Tech Stack")
content_slide(prs, "Technologies Used", [
    "Backend:      Python 3.9+, Flask, SQLite, Gunicorn",
    "Frontend:     HTML, CSS, Jinja2 templates",
    "AI:           OpenAI GPT-3.5-Turbo (optional)",
    "Messaging:    Twilio WhatsApp API (optional)",
    "Container:    Docker (multi-stage build)",
    "Orchestration: Docker Compose",
    "CI/CD:        GitHub Actions (lint → test → build)",
    "Quality:      Ruff linting, Pytest (27 tests)",
], two_col=True)

# ── Flow ───────────────────────────────────────────────────────────────
section_slide(prs, 6, "How It Works")
content_slide(prs, "Ticket Lifecycle", [
    "1. User describes issue via web form or WhatsApp",
    "2. Triage Engine analyzes text (free rules or OpenAI)",
    "3. Category + Priority assigned automatically",
    "   • Category: Hardware, Software, Network, Account, General",
    "   • Priority: Critical, High, Medium, Low",
    "4. Ticket saved to SQLite database",
    "5. User receives ticket ID and confirmation",
    "6. Admin views all tickets in dashboard",
    "7. Admin updates status: Open → In Progress → Resolved → Closed",
])

# ── DevOps ─────────────────────────────────────────────────────────────
section_slide(prs, 7, "DevOps & CI/CD")
content_slide(prs, "Pipeline", [
    "GitHub Actions (on push/PR to main):",
    "  1. Lint — ruff check . (imports, style, security)",
    "  2. Test — pytest -v (27 tests, isolated temp DB)",
    "  3. Build — docker build + docker compose build",
    "",
    "Docker:",
    "  • Multi-stage build (slim production image)",
    "  • Non-root user, HEALTHCHECK, Gunicorn server",
    "  • One-command deploy: docker compose up",
])

# ── Admin Dashboard ───────────────────────────────────────────────────
section_slide(prs, 8, "Admin Dashboard")
content_slide(prs, "Features", [
    "• Login at /admin (default password: admin)",
    "• View all tickets sorted by priority",
    "• Filter by status: Open, In Progress, Resolved, Closed",
    "• Update ticket status with one click",
    "• See full ticket details: user, description, category, priority",
])

# ── Triage ─────────────────────────────────────────────────────────────
section_slide(prs, 9, "Triage Engine")
content_slide(prs, "Free Rule-Based Mode", [
    "Categories (keyword matching):",
    "  Hardware: printer, keyboard, monitor, mouse, cable, broken, ...",
    "  Network: wifi, internet, connection, vpn, dns, offline, ...",
    "  Account: password, login, reset, 2fa, authentication, ...",
    "  Software: crash, error, install, update, bug, virus, ...",
    "  General: (fallback)",
    "",
    "Priorities:",
    "  Critical: down, outage, emergency, company-wide",
    "  High: important, asap, cannot, stuck, blocking",
    "  Medium: (default)   Low: question, minor, suggestion",
])

content_slide(prs, "OpenAI Mode (Optional)", [
    "When OPENAI_API_KEY is set:",
    "  • Uses GPT-3.5-Turbo for intelligent triage",
    "  • Prompt-engineered to return: Category | Priority",
    "  • Falls back to rule-based mode on API error",
    "",
    '  Example: "My computer won\'t boot up"',
    '  Response: "Hardware | Critical"',
])

# ── Project Structure ─────────────────────────────────────────────────
section_slide(prs, 10, "Project Structure")
content_slide(prs, "File Layout", [
    "├── bot.py              # Flask app (routes, admin, webhook)",
    "├── ai_triage.py        # Triage engine (rules + OpenAI)",
    "├── database.py         # SQLite database layer",
    "├── templates/          # 6 HTML templates",
    "├── tests/              # 27 pytest tests",
    "├── Dockerfile          # Multi-stage production build",
    "├── docker-compose.yml  # One-command deployment",
    "├── .github/workflows/  # CI/CD pipeline",
    "├── pyproject.toml      # Ruff + pytest config",
    "└── .env.example        # Environment variable template",
])

# ── Demo ───────────────────────────────────────────────────────────────
section_slide(prs, 11, "Demo")
content_slide(prs, "Quick Start", [
    "No API keys required — runs immediately:",
    "",
    "  pip install -r requirements.txt",
    "  python bot.py",
    "  → http://localhost:5000",
    "",
    "Or with Docker:",
    "",
    "  docker compose up",
    "  → http://localhost:5000",
    "",
    "Admin: /admin (password: admin)",
])

# ── Roadmap ────────────────────────────────────────────────────────────
section_slide(prs, 12, "Future Roadmap")
content_slide(prs, "Planned Features", [
    "• Email notifications on status changes",
    "• SMS alerts via Twilio",
    "• User authentication (citizen vs admin roles)",
    "• SLA tracking with auto-escalation",
    "• Reporting & analytics dashboard",
    "• File attachments (screenshots)",
    "• Multi-language support",
    "• Webhook API for third-party integration",
])

# ── Skills ─────────────────────────────────────────────────────────────
section_slide(prs, 13, "Skills Demonstrated")
content_slide(prs, "Engineering Competencies", [
    "Full-stack development: Flask + HTML/CSS + SQLite",
    "API integration: OpenAI GPT-3.5, Twilio WhatsApp",
    "DevOps: Docker, CI/CD, automated testing, linting",
    "AI/ML: GPT prompt engineering + rule-based engine",
    "UX design: Admin dashboard, responsive web forms",
    "Security: Session auth, non-root containers, env isolation",
    "Data: Relational database design, CRUD, query optimization",
])

# ── Thank You ──────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
text_box(slide, 1, 2, 11, 1, "Thank You", size=48, bold=True,
         color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)
text_box(slide, 1, 3.2, 11, 0.8, "Government ICT Service Ticket Bot",
         size=24, color=LIGHT_BLUE, alignment=PP_ALIGN.CENTER)
text_box(slide, 1, 4.5, 11, 0.5,
         "github.com/tamsirdev/gov-ict-ticket-bot",
         size=16, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)
text_box(slide, 1, 5.5, 11, 0.5,
         "Built with Python, Flask, Docker, and GitHub Actions",
         size=14, color=GRAY, alignment=PP_ALIGN.CENTER)


output = "ICT_Ticket_Bot_Presentation.pptx"
prs.save(output)
print(f"Presentation saved to {output}")
